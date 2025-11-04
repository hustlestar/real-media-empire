"""
Template manager for PowerPoint presentations.

Handles loading custom templates, creating default templates, and
managing style guidelines.
"""

import logging
import os
from pathlib import Path
from typing import Optional, Dict, Any, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .models import PresentationConfig, AspectRatio

logger = logging.getLogger(__name__)


# ============================================================================
# Layout Constants
# ============================================================================

# Standard layout indices in python-pptx
TITLE_LAYOUT = 0
CONTENT_LAYOUT = 1
SECTION_HEADER_LAYOUT = 2
TWO_COLUMN_LAYOUT = 3
COMPARISON_LAYOUT = 4
TITLE_ONLY_LAYOUT = 5
BLANK_LAYOUT = 6


# ============================================================================
# Template Manager
# ============================================================================


class TemplateManager:
    """
    Manages PowerPoint templates and style application.

    Supports:
    - Loading custom PPTX templates
    - Creating default templates from scratch
    - Extracting style guidelines from templates
    - Managing template placeholders
    """

    def __init__(self, config: PresentationConfig):
        """
        Initialize template manager.

        Args:
            config: Presentation configuration
        """
        self.config = config
        self.presentation: Optional[Presentation] = None
        self.style_guidelines: Dict[str, Any] = {}

    def load_or_create_presentation(self) -> Presentation:
        """
        Load template or create new presentation.

        Returns:
            Presentation object ready for content

        Raises:
            FileNotFoundError: If custom template path doesn't exist
        """
        if self.config.template_path:
            self.presentation = self.load_template(self.config.template_path)
            logger.info(f"Loaded template from {self.config.template_path}")
        else:
            self.presentation = self.create_default_template()
            logger.info("Created default template")

        self.style_guidelines = self.extract_style_guidelines()
        return self.presentation

    def load_template(self, template_path: str) -> Presentation:
        """
        Load existing PPTX template.

        Args:
            template_path: Path to .pptx template file

        Returns:
            Presentation object

        Raises:
            FileNotFoundError: If template doesn't exist
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")

        prs = Presentation(template_path)
        logger.info(f"Loaded template with {len(prs.slide_layouts)} layouts")
        return prs

    def create_default_template(self) -> Presentation:
        """
        Create a professional default template from scratch.

        Uses configuration settings for styling.

        Returns:
            New Presentation with default styling
        """
        prs = Presentation()

        # Set slide dimensions based on aspect ratio
        if self.config.aspect_ratio == AspectRatio.RATIO_16_9:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        else:  # 4:3
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

        # Apply theme colors
        self._apply_theme_colors(prs)

        logger.info(f"Created default {self.config.aspect_ratio.value} template")
        return prs

    def extract_style_guidelines(self) -> Dict[str, Any]:
        """
        Extract style guidelines from loaded template.

        Returns:
            Dictionary of style settings (fonts, colors, spacing)
        """
        if not self.presentation:
            return {}

        guidelines = {
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
            "num_layouts": len(self.presentation.slide_layouts),
            "theme": self.config.theme_name,
            "fonts": {"title": self.config.font_family, "body": self.config.font_family},
            "colors": {
                "primary": self.config.primary_color,
                "secondary": self.config.secondary_color,
                "accent": self.config.accent_color,
                "text": self.config.text_color,
                "background": self.config.background_color,
            },
        }

        return guidelines

    def get_layout(self, layout_name: str) -> Any:
        """
        Get slide layout by name.

        Args:
            layout_name: Name of layout ('title', 'content', etc.)

        Returns:
            Slide layout object
        """
        if not self.presentation:
            raise RuntimeError("No presentation loaded")

        layout_map = {
            "title": TITLE_LAYOUT,
            "content": CONTENT_LAYOUT,
            "section_header": SECTION_HEADER_LAYOUT,
            "two_column": TWO_COLUMN_LAYOUT,
            "blank": BLANK_LAYOUT,
        }

        layout_idx = layout_map.get(layout_name, CONTENT_LAYOUT)

        # Ensure layout exists
        if layout_idx >= len(self.presentation.slide_layouts):
            logger.warning(f"Layout {layout_idx} not found, using content layout")
            layout_idx = min(CONTENT_LAYOUT, len(self.presentation.slide_layouts) - 1)

        return self.presentation.slide_layouts[layout_idx]

    def save_presentation(self, output_path: str) -> str:
        """
        Save presentation to file.

        Args:
            output_path: Path to save .pptx file

        Returns:
            Absolute path to saved file

        Raises:
            RuntimeError: If no presentation loaded
        """
        if not self.presentation:
            raise RuntimeError("No presentation to save")

        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)

        self.presentation.save(output_path)
        abs_path = os.path.abspath(output_path)
        logger.info(f"Saved presentation to {abs_path}")
        return abs_path

    def _apply_theme_colors(self, prs: Presentation):
        """
        Apply theme colors to presentation.

        Args:
            prs: Presentation object to style
        """
        # Note: python-pptx has limited theme customization
        # Colors are primarily applied at slide/shape level
        # This method is a placeholder for future theme support
        pass

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """
        Convert hex color to RGB tuple.

        Args:
            hex_color: Hex color string (e.g., '#1F4E78')

        Returns:
            (r, g, b) tuple
        """
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


# ============================================================================
# Utility Functions
# ============================================================================


def get_default_template_path(aspect_ratio: AspectRatio) -> Optional[str]:
    """
    Get path to default template file if it exists.

    Args:
        aspect_ratio: Desired aspect ratio

    Returns:
        Path to template file, or None if not found
    """
    template_dir = Path(__file__).parent.parent.parent / "templates"

    if aspect_ratio == AspectRatio.RATIO_16_9:
        template_file = template_dir / "default_16x9.pptx"
    else:
        template_file = template_dir / "default_4x3.pptx"

    if template_file.exists():
        return str(template_file)

    return None
