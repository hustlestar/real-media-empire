"""
Slide builder for PowerPoint presentations.

Wrapper around python-pptx for creating and formatting slides.
"""

import logging
import re
from typing import Optional, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.dml.color import RGBColor

from .models import SlideDefinition, SlideLayout, BulletPoint, PresentationConfig
from .template_manager import TemplateManager

logger = logging.getLogger(__name__)


# ============================================================================
# Slide Builder
# ============================================================================


class SlideBuilder:
    """
    Builds PowerPoint slides from SlideDefinition models.

    Handles:
    - Creating slides with header, body, and footer sections
    - Adding and formatting text
    - Managing bullet points and indentation
    - Auto-fitting text to shapes
    """

    def __init__(self, template_manager: TemplateManager, config: PresentationConfig):
        """
        Initialize slide builder.

        Args:
            template_manager: Template manager with loaded presentation
            config: Presentation configuration
        """
        self.template_manager = template_manager
        self.config = config
        self.presentation = template_manager.presentation

        if not self.presentation:
            raise RuntimeError("Template manager has no presentation loaded")

        # Get dimensions from config
        self.HEADER_HEIGHT = config.header_height
        self.FOOTER_HEIGHT = config.footer_height
        self.MARGIN = config.margin

        # Calculate body dimensions based on slide size
        self.slide_width = self.presentation.slide_width.inches
        self.slide_height = self.presentation.slide_height.inches
        self.body_top = self.HEADER_HEIGHT
        self.body_height = self.slide_height - self.HEADER_HEIGHT - self.FOOTER_HEIGHT
        self.footer_top = self.slide_height - self.FOOTER_HEIGHT

    def build_slide(self, slide_def: SlideDefinition) -> None:
        """
        Build a slide from definition.

        Args:
            slide_def: Slide definition with content and layout
        """
        layout_name = self._map_layout_to_name(slide_def.layout_type)

        if slide_def.layout_type == SlideLayout.TITLE:
            self.create_title_slide(slide_def)
        elif slide_def.layout_type == SlideLayout.TWO_COLUMN:
            self.create_two_column_slide(slide_def)
        elif slide_def.layout_type == SlideLayout.IMAGE_TEXT:
            self.create_image_text_slide(slide_def)
        elif slide_def.layout_type == SlideLayout.BLANK:
            self.create_blank_slide(slide_def)
        else:
            self.create_content_slide(slide_def)

        logger.debug(f"Built slide {slide_def.slide_number}: {slide_def.content.title}")

    def create_title_slide(self, slide_def: SlideDefinition) -> None:
        """
        Create title slide with full-screen header and centered content.

        Args:
            slide_def: Slide definition
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Use blank layout for full control
        blank_layout = self.template_manager.get_layout("blank")
        slide = self.presentation.slides.add_slide(blank_layout)

        # Create full-height header (no body/footer on title slide)
        title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(self.slide_width), Inches(self.slide_height))

        # Style background with gradient-like effect (primary color)
        title_bg.fill.solid()
        bg_color = self._hex_to_rgb(self.config.primary_color)
        title_bg.fill.fore_color.rgb = RGBColor(*bg_color)
        title_bg.line.fill.background()

        # Add title text (centered)
        title_box = slide.shapes.add_textbox(
            Inches(self.MARGIN), Inches(self.slide_height / 2 - 1.5), Inches(self.slide_width - 2 * self.MARGIN), Inches(1.5)
        )

        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Use formatted text parser for title
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # Parse and add formatted title text
        segments = self._parse_formatted_text(slide_def.content.title)
        for idx, (segment_text, is_bold, is_italic) in enumerate(segments):
            if idx == 0:
                p.text = segment_text
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = segment_text

            run.font.size = Pt(44)
            run.font.bold = True  # Title is always bold
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.name = self.config.font_family

        # Add subtitle if present
        subtitle_text = ""
        if slide_def.content.text_content:
            subtitle_text = slide_def.content.text_content
        elif slide_def.content.bullets:
            subtitle_text = slide_def.content.bullets[0].text

        if subtitle_text:
            subtitle_box = slide.shapes.add_textbox(
                Inches(self.MARGIN), Inches(self.slide_height / 2 + 0.5), Inches(self.slide_width - 2 * self.MARGIN), Inches(1.0)
            )

            text_frame = subtitle_box.text_frame
            text_frame.word_wrap = True

            # Use formatted text parser for subtitle
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            # Parse and add formatted subtitle text
            segments = self._parse_formatted_text(subtitle_text)
            for idx, (segment_text, is_bold, is_italic) in enumerate(segments):
                if idx == 0:
                    p.text = segment_text
                    run = p.runs[0]
                else:
                    run = p.add_run()
                    run.text = segment_text

                run.font.size = Pt(24)
                run.font.bold = is_bold
                run.font.italic = is_italic
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.name = self.config.font_family

    def create_content_slide(self, slide_def: SlideDefinition) -> None:
        """
        Create standard content slide with header, body, and footer.

        Args:
            slide_def: Slide definition
        """
        # Use blank layout to have full control
        blank_layout = self.template_manager.get_layout("blank")
        slide = self.presentation.slides.add_slide(blank_layout)

        # Add header section
        self._add_header(slide, slide_def.content.title)

        # Add body section with content
        self._add_body_with_bullets(slide, slide_def.content.bullets)

        # Add footer section
        self._add_footer(slide, slide_def.slide_number)

        # Add speaker notes if enabled
        if self.config.enable_speaker_notes and slide_def.content.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_def.content.speaker_notes

    def create_two_column_slide(self, slide_def: SlideDefinition) -> None:
        """
        Create two-column layout slide.

        Args:
            slide_def: Slide definition
        """
        # Try to get two-column layout, fall back to content
        try:
            layout = self.template_manager.get_layout("two_column")
        except:
            layout = self.template_manager.get_layout("content")

        slide = self.presentation.slides.add_slide(layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_def.content.title
        self._apply_text_styling(title.text_frame, is_title=True)

        # Split bullets between columns
        if slide_def.content.bullets:
            mid_point = len(slide_def.content.bullets) // 2
            left_bullets = slide_def.content.bullets[:mid_point]
            right_bullets = slide_def.content.bullets[mid_point:]

            # This is simplified - actual two-column would need more layout work
            # For now, just add all bullets to content area
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    self._add_bullet_points(text_frame, slide_def.content.bullets)
                    break

    def create_image_text_slide(self, slide_def: SlideDefinition) -> None:
        """
        Create slide with image placeholder and text.

        Args:
            slide_def: Slide definition
        """
        # For now, create as content slide
        # Image integration would require actual image files
        self.create_content_slide(slide_def)

    def create_blank_slide(self, slide_def: SlideDefinition) -> None:
        """
        Create blank slide.

        Args:
            slide_def: Slide definition
        """
        layout = self.template_manager.get_layout("blank")
        slide = self.presentation.slides.add_slide(layout)

        # Add title as text box if present
        if slide_def.content.title:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = slide_def.content.title

            self._apply_text_styling(text_frame, is_title=True)

    def _parse_formatted_text(self, text: str) -> List[Tuple[str, bool, bool]]:
        """
        Parse markdown-style formatting from text.

        Args:
            text: Text with **bold** and *italic* markers

        Returns:
            List of (text, is_bold, is_italic) tuples
        """
        segments = []
        i = 0

        while i < len(text):
            # Check for **bold**
            if i < len(text) - 3 and text[i : i + 2] == "**":
                # Find closing **
                end = text.find("**", i + 2)
                if end != -1:
                    # Found bold text
                    bold_text = text[i + 2 : end]
                    if bold_text:  # Only add if not empty
                        segments.append((bold_text, True, False))
                    i = end + 2
                    continue

            # Check for *italic* (but not **)
            if i < len(text) - 2 and text[i] == "*" and text[i + 1] != "*":
                # Find closing *
                end = text.find("*", i + 1)
                if end != -1 and (end == len(text) - 1 or text[end + 1] != "*"):
                    # Found italic text
                    italic_text = text[i + 1 : end]
                    if italic_text:  # Only add if not empty
                        segments.append((italic_text, False, True))
                    i = end + 1
                    continue

            # Regular character - accumulate plain text
            plain_start = i
            while i < len(text) and text[i] != "*":
                i += 1

            if i > plain_start:
                segments.append((text[plain_start:i], False, False))

        # If no segments were found, return the original text as plain
        if not segments:
            segments.append((text, False, False))

        return segments

    def _add_formatted_text_to_paragraph(self, paragraph, text: str, base_font_size: int):
        """
        Add formatted text to paragraph with bold/italic support.

        Args:
            paragraph: Paragraph object to add text to
            text: Text with markdown formatting
            base_font_size: Base font size in points
        """
        segments = self._parse_formatted_text(text)

        for idx, (segment_text, is_bold, is_italic) in enumerate(segments):
            if idx == 0:
                # First segment uses the existing paragraph run
                paragraph.text = segment_text
                run = paragraph.runs[0]
            else:
                # Subsequent segments add new runs
                run = paragraph.add_run()
                run.text = segment_text

            # Apply formatting
            run.font.size = Pt(base_font_size)
            run.font.name = self.config.font_family
            run.font.bold = is_bold
            run.font.italic = is_italic

            # Apply text color
            if self.config.text_color:
                rgb = self._hex_to_rgb(self.config.text_color)
                run.font.color.rgb = RGBColor(*rgb)

    def _add_bullet_points(self, text_frame, bullets: List[BulletPoint]) -> None:
        """
        Add bullet points to text frame with proper indentation and formatting.

        Args:
            text_frame: Text frame to add bullets to
            bullets: List of bullet points
        """
        from .models import ListType

        # Track numbering for each level
        number_counters = {0: 0, 1: 0, 2: 0}
        last_list_type = None

        for idx, bullet in enumerate(bullets):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.level = bullet.level

            # Calculate font size based on level with better hierarchy
            if bullet.level == 0:
                base_font_size = self.config.body_font_size
            elif bullet.level == 1:
                base_font_size = self.config.body_font_size - 4
            else:
                base_font_size = self.config.body_font_size - 6

            # Handle numbering
            if bullet.list_type == ListType.NUMBERED:
                # Reset counter if switching from bullets to numbers at same level
                if last_list_type != ListType.NUMBERED or (idx > 0 and bullets[idx - 1].level != bullet.level):
                    if idx > 0 and bullets[idx - 1].level > bullet.level:
                        # Reset counter when going back to parent level
                        number_counters[bullet.level] = 0

                # Increment counter for this level
                number_counters[bullet.level] += 1

                # Prepend number to text
                number_prefix = f"{number_counters[bullet.level]}. "
                formatted_text = number_prefix + bullet.text
            else:
                formatted_text = bullet.text
                # Reset number counter if we switch to bullets
                if last_list_type == ListType.NUMBERED:
                    number_counters = {0: 0, 1: 0, 2: 0}

            last_list_type = bullet.list_type

            # Add spacing between bullet points for readability
            if bullet.level == 0 and idx > 0:
                p.space_before = Pt(12)  # Add space before main bullets
            else:
                p.space_before = Pt(6)  # Less space for sub-bullets

            # Set line spacing
            p.line_spacing = 1.2  # 120% line spacing

            # Left alignment
            p.alignment = PP_ALIGN.LEFT

            # Add formatted text to paragraph
            self._add_formatted_text_to_paragraph(p, formatted_text, base_font_size)

    def _apply_text_styling(self, text_frame, is_title: bool = False) -> None:
        """
        Apply text styling to text frame.

        Args:
            text_frame: Text frame to style
            is_title: Whether this is a title (larger font)
        """
        for paragraph in text_frame.paragraphs:
            if is_title:
                paragraph.font.size = Pt(self.config.title_font_size)
                paragraph.font.bold = True
            else:
                paragraph.font.size = Pt(self.config.body_font_size)

            paragraph.font.name = self.config.font_family

            # Apply color if specified
            if self.config.text_color:
                rgb = self._hex_to_rgb(self.config.text_color)
                paragraph.font.color.rgb = RGBColor(*rgb)

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def _map_layout_to_name(self, layout_type: SlideLayout) -> str:
        """Map SlideLayout enum to layout name."""
        mapping = {
            SlideLayout.TITLE: "title",
            SlideLayout.CONTENT: "content",
            SlideLayout.TWO_COLUMN: "two_column",
            SlideLayout.IMAGE_TEXT: "content",  # Fallback
            SlideLayout.BLANK: "blank",
            SlideLayout.SECTION_HEADER: "section_header",
        }
        return mapping.get(layout_type, "content")

    def fit_text_to_shape(self, text_frame, max_font_size: int = 32, min_font_size: int = 10) -> None:
        """
        Auto-fit text to shape by reducing font size if needed.

        Args:
            text_frame: Text frame to fit
            max_font_size: Maximum font size
            min_font_size: Minimum font size
        """
        # Note: python-pptx has limited auto-fit support
        # This is a placeholder for more sophisticated fitting
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Don't auto-resize shape

    # ========================================================================
    # Header, Body, Footer Methods
    # ========================================================================

    def _add_header(self, slide, title_text: str) -> None:
        """
        Add header section to slide with title.

        Args:
            slide: Slide object
            title_text: Title text for the header
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Create accent line at top (thin colored bar)
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(self.slide_width), Inches(0.08))  # Thin accent line
        accent_line.fill.solid()
        accent_color = self._hex_to_rgb(self.config.primary_color)
        accent_line.fill.fore_color.rgb = RGBColor(*accent_color)
        accent_line.line.fill.background()

        # Add title as text box (not a full background rectangle)
        title_box = slide.shapes.add_textbox(
            Inches(self.MARGIN * 1.5), Inches(0.15), Inches(self.slide_width - 3 * self.MARGIN), Inches(self.HEADER_HEIGHT - 0.15)
        )

        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Style title text - larger and bold
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        # Use primary color for title text
        title_color = self._hex_to_rgb(self.config.primary_color)

        # Parse and add formatted title text
        segments = self._parse_formatted_text(title_text)
        for idx, (segment_text, is_bold, is_italic) in enumerate(segments):
            if idx == 0:
                p.text = segment_text
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = segment_text

            run.font.size = Pt(36)
            run.font.bold = True  # Header titles are always bold
            run.font.name = self.config.font_family
            run.font.color.rgb = RGBColor(*title_color)

    def _add_body_with_bullets(self, slide, bullets: List[BulletPoint]) -> None:
        """
        Add body section with bullet points.

        Args:
            slide: Slide object
            bullets: List of bullet points
        """
        from pptx.enum.shapes import MSO_SHAPE

        # Create body rectangle
        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(self.body_top), Inches(self.slide_width), Inches(self.body_height))

        # Style body background (white)
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(255, 255, 255)
        body.line.fill.background()  # No border

        # Add bullet points
        text_frame = body.text_frame
        text_frame.clear()

        # Better margins for content
        text_frame.margin_left = Inches(self.MARGIN * 1.5)
        text_frame.margin_right = Inches(self.MARGIN * 1.5)
        text_frame.margin_top = Inches(self.MARGIN * 1.2)
        text_frame.margin_bottom = Inches(self.MARGIN)
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        if bullets:
            self._add_bullet_points(text_frame, bullets)
        else:
            # Add placeholder text if no bullets
            text_frame.text = "[Content area]"

    def _add_footer(self, slide, slide_number: int) -> None:
        """
        Add footer section to slide.

        Args:
            slide: Slide object
            slide_number: Slide number to display
        """
        from pptx.enum.shapes import MSO_SHAPE
        from datetime import datetime

        # Create footer rectangle
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(self.footer_top), Inches(self.slide_width), Inches(self.FOOTER_HEIGHT))

        # Style footer background (light gray)
        footer.fill.solid()
        footer_color = self._hex_to_rgb("#E7E7E7")
        footer.fill.fore_color.rgb = RGBColor(*footer_color)
        footer.line.fill.background()  # No border

        # Add footer text
        text_frame = footer.text_frame
        text_frame.clear()

        # Format: "Slide X | Date"
        footer_text = f"Slide {slide_number} | {datetime.now().strftime('%B %Y')}"
        text_frame.text = footer_text
        text_frame.margin_top = Inches(0.15)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Style footer text
        p = text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(102, 102, 102)  # Gray text
        p.font.name = self.config.font_family
        p.alignment = PP_ALIGN.CENTER
