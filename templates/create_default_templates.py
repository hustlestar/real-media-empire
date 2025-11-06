"""
Script to create default PowerPoint templates with header, body, and footer.

Run this to generate default 4:3 and 16:9 templates.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_template_with_layouts(width, height, filename, aspect_ratio_name):
    """
    Create a template with proper header, body, and footer structure.

    Args:
        width: Slide width in inches
        height: Slide height in inches
        filename: Output filename
        aspect_ratio_name: Name for display (e.g., "16:9")
    """
    prs = Presentation()

    # Set dimensions
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    # Colors
    HEADER_BG = hex_to_rgb('#1F4E78')  # Dark blue
    FOOTER_BG = hex_to_rgb('#E7E7E7')  # Light gray
    ACCENT = hex_to_rgb('#2E75B6')     # Medium blue
    TEXT_DARK = hex_to_rgb('#000000')  # Black
    TEXT_LIGHT = hex_to_rgb('#FFFFFF') # White

    # Dimensions (adjusted for aspect ratio)
    margin = 0.5
    header_height = 0.8
    footer_height = 0.6

    content_top = header_height + margin
    content_bottom = height - footer_height - margin
    content_height = content_bottom - content_top

    # Layout 0: Title Slide
    title_layout = prs.slide_layouts[0]

    # Layout 1: Title and Content (with header, body, footer)
    # This is the main layout we'll modify
    content_layout = prs.slide_layouts[1]

    print(f"Created template: {filename}")
    print(f"  - Aspect ratio: {aspect_ratio_name}")
    print(f"  - Dimensions: {width}\" x {height}\"")
    print(f"  - Header height: {header_height}\"")
    print(f"  - Footer height: {footer_height}\"")
    print(f"  - Content area: {content_height}\"")

    # Save
    prs.save(filename)
    print(f"  - Saved successfully\n")


def create_professional_template(width, height, filename, aspect_ratio_name):
    """
    Create a professional template from scratch with explicit header/body/footer.

    This creates a blank presentation and adds a slide with shapes for header, body, footer.
    """
    prs = Presentation()

    # Set dimensions
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    # Colors
    HEADER_BG = hex_to_rgb('#1F4E78')  # Dark blue
    FOOTER_BG = hex_to_rgb('#E7E7E7')  # Light gray
    ACCENT = hex_to_rgb('#2E75B6')     # Medium blue
    BODY_BG = hex_to_rgb('#FFFFFF')    # White

    # Dimensions
    margin = 0.5
    header_height = 1.0
    footer_height = 0.5

    # Create a blank slide to demonstrate the structure
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add Header Rectangle
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        Inches(0),
        Inches(width),
        Inches(header_height)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*HEADER_BG)
    header.line.color.rgb = RGBColor(*HEADER_BG)

    # Add header text
    header_text = header.text_frame
    header_text.text = "Presentation Title"
    header_text.paragraphs[0].font.size = Pt(28)
    header_text.paragraphs[0].font.bold = True
    header_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header_text.paragraphs[0].alignment = PP_ALIGN.LEFT
    header_text.margin_left = Inches(margin)
    header_text.margin_top = Inches(0.25)

    # Add Body Rectangle
    body_top = header_height
    body_height = height - header_height - footer_height

    body = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        Inches(body_top),
        Inches(width),
        Inches(body_height)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(*BODY_BG)
    body.line.width = Pt(0)  # No border

    # Add body content placeholder text
    body_text = body.text_frame
    body_text.text = "Content Area\n• Bullet point 1\n• Bullet point 2\n• Bullet point 3"
    body_text.paragraphs[0].font.size = Pt(24)
    body_text.paragraphs[0].font.bold = True
    body_text.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb('#000000'))
    body_text.margin_left = Inches(margin)
    body_text.margin_right = Inches(margin)
    body_text.margin_top = Inches(margin)
    body_text.word_wrap = True

    # Add Footer Rectangle
    footer_top = height - footer_height

    footer = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0),
        Inches(footer_top),
        Inches(width),
        Inches(footer_height)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(*FOOTER_BG)
    footer.line.color.rgb = RGBColor(*FOOTER_BG)

    # Add footer text
    footer_text = footer.text_frame
    footer_text.text = "Footer Text | Slide Number | Date"
    footer_text.paragraphs[0].font.size = Pt(12)
    footer_text.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb('#666666'))
    footer_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_text.margin_top = Inches(0.15)

    print(f"Created professional template: {filename}")
    print(f"  - Aspect ratio: {aspect_ratio_name}")
    print(f"  - Dimensions: {width}\" x {height}\"")
    print(f"  - Header: 0\" to {header_height}\" (dark blue background)")
    print(f"  - Body: {header_height}\" to {footer_top}\" (white background)")
    print(f"  - Footer: {footer_top}\" to {height}\" (light gray background)")

    # Save
    prs.save(filename)
    print(f"  - Saved successfully\n")


def create_16x9_template():
    """Create default 16:9 template."""
    create_professional_template(
        width=13.333,
        height=7.5,
        filename='default_16x9.pptx',
        aspect_ratio_name='16:9'
    )


def create_4x3_template():
    """Create default 4:3 template."""
    create_professional_template(
        width=10,
        height=7.5,
        filename='default_4x3.pptx',
        aspect_ratio_name='4:3'
    )


if __name__ == '__main__':
    print("="*60)
    print("Creating Professional PPTX Templates")
    print("with Header, Body, and Footer Structure")
    print("="*60)
    print()

    create_16x9_template()
    create_4x3_template()

    print("="*60)
    print("Done! Templates created in current directory.")
    print("="*60)
    print()
    print("Template Structure:")
    print("  [HEADER]  - Company/presentation branding area")
    print("  [BODY]    - Main content area")
    print("  [FOOTER]  - Page numbers, date, additional info")
